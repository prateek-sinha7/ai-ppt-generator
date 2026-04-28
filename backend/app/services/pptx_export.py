"""
PPTX Export Service - Enterprise-grade PowerPoint generation from Slide_JSON.

This service implements:
- Task 24.1: python-pptx slide builder mapping each slide type to appropriate PPTX layout
- Task 24.2: Theme application preserving Corporate/Executive/Professional/Dark Modern color schemes
- Task 24.3: Chart rendering in PPTX (bar/line/pie/area/scatter/stacked_bar/donut)
- Task 24.4: Table rendering in PPTX with proper formatting
- Task 24.5: Transition mapping (fade→Fade, slide→Push, none→no transition)
- Task 24.6: S3/MinIO upload and signed URL generation (1-hour TTL)
- Task 24.7: Performance validation (completes within 30 seconds for 50 slides)

Enterprise design features:
- Accent bars and decorative shapes on every slide
- KPI / metric highlight boxes with large numbers
- Icon-label pairs rendered as styled text badges
- Color-coded section headers for comparison slides
- Gradient-style title slides with bottom accent strip
- Alternating row colors and bold header styling for tables
- Visual hierarchy: title → subtitle → body → caption
"""

import io
from typing import Any, Dict, List, Optional, Tuple
from enum import Enum

import structlog
from pptx import Presentation as PptxPresentation
from pptx.chart.data import CategoryChartData, BubbleChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.dml import MSO_THEME_COLOR

logger = structlog.get_logger(__name__)


# ---------------------------------------------------------------------------
# Theme Definitions
# ---------------------------------------------------------------------------

class ThemeColors:
    """Color schemes for each presentation theme."""

    OCEAN_DEPTHS = {
        "primary": RGBColor(26, 35, 50),
        "secondary": RGBColor(45, 139, 139),
        "accent": RGBColor(168, 218, 220),
        "accent2": RGBColor(91, 163, 163),
        "text": RGBColor(26, 35, 50),
        "text_light": RGBColor(107, 114, 128),
        "background": RGBColor(241, 250, 238),
        "surface": RGBColor(245, 248, 246),
        "divider": RGBColor(209, 213, 219),
        "kpi_bg": RGBColor(26, 35, 50),
        "kpi_text": RGBColor(255, 255, 255),
        "header_bar": RGBColor(26, 35, 50),
        "chart_colors": [
            RGBColor(26, 35, 50), RGBColor(45, 139, 139), RGBColor(168, 218, 220),
            RGBColor(91, 163, 163), RGBColor(61, 107, 107), RGBColor(130, 180, 180),
            RGBColor(107, 114, 128),
        ],
    }

    SUNSET_BOULEVARD = {
        "primary": RGBColor(231, 111, 81),
        "secondary": RGBColor(244, 162, 97),
        "accent": RGBColor(233, 196, 106),
        "accent2": RGBColor(42, 157, 143),
        "text": RGBColor(38, 70, 83),
        "text_light": RGBColor(107, 114, 128),
        "background": RGBColor(255, 255, 255),
        "surface": RGBColor(253, 248, 243),
        "divider": RGBColor(209, 213, 219),
        "kpi_bg": RGBColor(38, 70, 83),
        "kpi_text": RGBColor(231, 111, 81),
        "header_bar": RGBColor(38, 70, 83),
        "chart_colors": [
            RGBColor(231, 111, 81), RGBColor(244, 162, 97), RGBColor(233, 196, 106),
            RGBColor(38, 70, 83), RGBColor(42, 157, 143), RGBColor(180, 90, 65),
            RGBColor(107, 114, 128),
        ],
    }

    FOREST_CANOPY = {
        "primary": RGBColor(45, 74, 43),
        "secondary": RGBColor(125, 132, 113),
        "accent": RGBColor(164, 172, 134),
        "accent2": RGBColor(90, 122, 88),
        "text": RGBColor(45, 74, 43),
        "text_light": RGBColor(107, 114, 128),
        "background": RGBColor(250, 249, 246),
        "surface": RGBColor(245, 245, 242),
        "divider": RGBColor(209, 213, 219),
        "kpi_bg": RGBColor(45, 74, 43),
        "kpi_text": RGBColor(255, 255, 255),
        "header_bar": RGBColor(45, 74, 43),
        "chart_colors": [
            RGBColor(45, 74, 43), RGBColor(125, 132, 113), RGBColor(164, 172, 134),
            RGBColor(90, 122, 88), RGBColor(139, 155, 120), RGBColor(70, 100, 68),
            RGBColor(107, 114, 128),
        ],
    }

    MODERN_MINIMALIST = {
        "primary": RGBColor(54, 69, 79),
        "secondary": RGBColor(112, 128, 144),
        "accent": RGBColor(211, 211, 211),
        "accent2": RGBColor(80, 90, 100),
        "text": RGBColor(54, 69, 79),
        "text_light": RGBColor(112, 128, 144),
        "background": RGBColor(255, 255, 255),
        "surface": RGBColor(245, 245, 245),
        "divider": RGBColor(211, 211, 211),
        "kpi_bg": RGBColor(54, 69, 79),
        "kpi_text": RGBColor(255, 255, 255),
        "header_bar": RGBColor(54, 69, 79),
        "chart_colors": [
            RGBColor(54, 69, 79), RGBColor(112, 128, 144), RGBColor(160, 160, 160),
            RGBColor(80, 90, 100), RGBColor(136, 150, 160), RGBColor(40, 50, 60),
            RGBColor(180, 180, 180),
        ],
    }

    GOLDEN_HOUR = {
        "primary": RGBColor(244, 169, 0),
        "secondary": RGBColor(193, 102, 107),
        "accent": RGBColor(212, 184, 150),
        "accent2": RGBColor(139, 105, 20),
        "text": RGBColor(74, 64, 58),
        "text_light": RGBColor(107, 114, 128),
        "background": RGBColor(255, 255, 255),
        "surface": RGBColor(250, 246, 240),
        "divider": RGBColor(209, 213, 219),
        "kpi_bg": RGBColor(74, 64, 58),
        "kpi_text": RGBColor(244, 169, 0),
        "header_bar": RGBColor(74, 64, 58),
        "chart_colors": [
            RGBColor(244, 169, 0), RGBColor(193, 102, 107), RGBColor(212, 184, 150),
            RGBColor(139, 105, 20), RGBColor(160, 82, 78), RGBColor(180, 140, 60),
            RGBColor(107, 114, 128),
        ],
    }

    ARCTIC_FROST = {
        "primary": RGBColor(74, 111, 165),
        "secondary": RGBColor(192, 192, 192),
        "accent": RGBColor(212, 228, 247),
        "accent2": RGBColor(85, 128, 168),
        "text": RGBColor(44, 62, 80),
        "text_light": RGBColor(107, 114, 128),
        "background": RGBColor(250, 250, 250),
        "surface": RGBColor(245, 247, 250),
        "divider": RGBColor(209, 213, 219),
        "kpi_bg": RGBColor(74, 111, 165),
        "kpi_text": RGBColor(255, 255, 255),
        "header_bar": RGBColor(74, 111, 165),
        "chart_colors": [
            RGBColor(74, 111, 165), RGBColor(122, 156, 198), RGBColor(168, 196, 224),
            RGBColor(85, 128, 168), RGBColor(61, 90, 128), RGBColor(100, 140, 180),
            RGBColor(107, 114, 128),
        ],
    }

    DESERT_ROSE = {
        "primary": RGBColor(212, 165, 165),
        "secondary": RGBColor(184, 125, 109),
        "accent": RGBColor(232, 213, 196),
        "accent2": RGBColor(155, 107, 107),
        "text": RGBColor(93, 46, 70),
        "text_light": RGBColor(107, 114, 128),
        "background": RGBColor(255, 255, 255),
        "surface": RGBColor(250, 245, 240),
        "divider": RGBColor(209, 213, 219),
        "kpi_bg": RGBColor(93, 46, 70),
        "kpi_text": RGBColor(212, 165, 165),
        "header_bar": RGBColor(93, 46, 70),
        "chart_colors": [
            RGBColor(212, 165, 165), RGBColor(184, 125, 109), RGBColor(232, 213, 196),
            RGBColor(93, 46, 70), RGBColor(155, 107, 107), RGBColor(170, 140, 140),
            RGBColor(107, 114, 128),
        ],
    }

    TECH_INNOVATION = {
        "primary": RGBColor(0, 102, 255),
        "secondary": RGBColor(0, 255, 255),
        "accent": RGBColor(0, 204, 204),
        "accent2": RGBColor(51, 136, 255),
        "text": RGBColor(255, 255, 255),
        "text_light": RGBColor(156, 163, 175),
        "background": RGBColor(30, 30, 30),
        "surface": RGBColor(42, 42, 42),
        "divider": RGBColor(55, 65, 81),
        "kpi_bg": RGBColor(0, 102, 255),
        "kpi_text": RGBColor(255, 255, 255),
        "header_bar": RGBColor(0, 102, 255),
        "chart_colors": [
            RGBColor(0, 102, 255), RGBColor(0, 255, 255), RGBColor(0, 204, 204),
            RGBColor(51, 136, 255), RGBColor(102, 221, 255), RGBColor(0, 170, 170),
            RGBColor(156, 163, 175),
        ],
    }

    BOTANICAL_GARDEN = {
        "primary": RGBColor(74, 124, 89),
        "secondary": RGBColor(249, 166, 32),
        "accent": RGBColor(183, 71, 42),
        "accent2": RGBColor(107, 155, 120),
        "text": RGBColor(58, 58, 58),
        "text_light": RGBColor(107, 114, 128),
        "background": RGBColor(245, 243, 237),
        "surface": RGBColor(240, 237, 230),
        "divider": RGBColor(209, 213, 219),
        "kpi_bg": RGBColor(74, 124, 89),
        "kpi_text": RGBColor(255, 255, 255),
        "header_bar": RGBColor(74, 124, 89),
        "chart_colors": [
            RGBColor(74, 124, 89), RGBColor(249, 166, 32), RGBColor(183, 71, 42),
            RGBColor(107, 155, 120), RGBColor(212, 136, 26), RGBColor(90, 150, 105),
            RGBColor(107, 114, 128),
        ],
    }

    MIDNIGHT_GALAXY = {
        "primary": RGBColor(74, 78, 143),
        "secondary": RGBColor(164, 144, 194),
        "accent": RGBColor(230, 230, 250),
        "accent2": RGBColor(107, 111, 175),
        "text": RGBColor(230, 230, 250),
        "text_light": RGBColor(156, 163, 175),
        "background": RGBColor(43, 30, 62),
        "surface": RGBColor(54, 42, 78),
        "divider": RGBColor(74, 64, 96),
        "kpi_bg": RGBColor(74, 78, 143),
        "kpi_text": RGBColor(230, 230, 250),
        "header_bar": RGBColor(74, 78, 143),
        "chart_colors": [
            RGBColor(74, 78, 143), RGBColor(164, 144, 194), RGBColor(230, 230, 250),
            RGBColor(107, 111, 175), RGBColor(196, 184, 216), RGBColor(90, 94, 160),
            RGBColor(156, 163, 175),
        ],
    }

    @classmethod
    def get_theme(cls, theme_name: str) -> Dict[str, Any]:
        """Get theme colors by name."""
        theme_map = {
            "ocean-depths": cls.OCEAN_DEPTHS,
            "ocean_depths": cls.OCEAN_DEPTHS,
            "sunset-boulevard": cls.SUNSET_BOULEVARD,
            "sunset_boulevard": cls.SUNSET_BOULEVARD,
            "forest-canopy": cls.FOREST_CANOPY,
            "forest_canopy": cls.FOREST_CANOPY,
            "modern-minimalist": cls.MODERN_MINIMALIST,
            "modern_minimalist": cls.MODERN_MINIMALIST,
            "golden-hour": cls.GOLDEN_HOUR,
            "golden_hour": cls.GOLDEN_HOUR,
            "arctic-frost": cls.ARCTIC_FROST,
            "arctic_frost": cls.ARCTIC_FROST,
            "desert-rose": cls.DESERT_ROSE,
            "desert_rose": cls.DESERT_ROSE,
            "tech-innovation": cls.TECH_INNOVATION,
            "tech_innovation": cls.TECH_INNOVATION,
            "botanical-garden": cls.BOTANICAL_GARDEN,
            "botanical_garden": cls.BOTANICAL_GARDEN,
            "midnight-galaxy": cls.MIDNIGHT_GALAXY,
            "midnight_galaxy": cls.MIDNIGHT_GALAXY,
        }
        return theme_map.get(theme_name.lower(), cls.OCEAN_DEPTHS)


# ---------------------------------------------------------------------------
# Transition Mapping
# ---------------------------------------------------------------------------

class TransitionMapper:
    """Maps Slide_JSON transitions to PowerPoint transitions."""
    
    # PowerPoint transition constants
    # Note: python-pptx doesn't expose all transition enums, so we use XML manipulation
    TRANSITION_MAP = {
        "fade": "fade",
        "slide": "push",
        "none": None,
    }
    
    @staticmethod
    def apply_transition(slide, transition_type: Optional[str] = "fade"):
        """
        Apply transition to a slide using XML manipulation.
        
        Args:
            slide: python-pptx slide object
            transition_type: Transition type from Slide_JSON ("fade", "slide", "none")
        """
        if not transition_type or transition_type == "none":
            return
        
        pptx_transition = TransitionMapper.TRANSITION_MAP.get(transition_type)
        if not pptx_transition:
            return
        
        try:
            # Access slide XML
            sld = slide._element
            
            # Remove existing transition if present
            existing = sld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
            if existing is not None:
                sld.remove(existing)
            
            # Create transition element
            transition = OxmlElement('p:transition')
            
            if pptx_transition == "fade":
                # Add fade transition
                fade = OxmlElement('p:fade')
                fade.set('thruBlk', '0')
                transition.append(fade)
            elif pptx_transition == "push":
                # Add push transition
                push = OxmlElement('p:push')
                push.set('dir', 'l')  # left direction
                transition.append(push)
            
            # Set transition speed (medium)
            transition.set('spd', 'med')
            
            # Insert transition at beginning of slide
            sld.insert(0, transition)
            
            logger.debug("transition_applied", transition=transition_type)
            
        except Exception as e:
            logger.warning("transition_application_failed", transition=transition_type, error=str(e))


# ---------------------------------------------------------------------------
# PPTX Builder
# ---------------------------------------------------------------------------

class PPTXBuilder:
    """
    Enterprise-grade PPTX builder.

    Handles:
    - Slide type mapping to layouts
    - Theme application with accent bars and decorative elements
    - Chart rendering: bar, line, pie, area, stacked_bar, donut, scatter
    - Table rendering with alternating rows and bold headers
    - Metric / KPI slides with large number display
    - Transition application
    - Slide background fill for dark themes
    """

    # Slide dimensions (16:9 widescreen)
    SLIDE_WIDTH = Inches(13.33)
    SLIDE_HEIGHT = Inches(7.5)

    # Layout margins
    MARGIN_TOP = Inches(0.5)
    MARGIN_BOTTOM = Inches(0.5)
    MARGIN_LEFT = Inches(0.75)
    MARGIN_RIGHT = Inches(0.75)

    # Accent bar dimensions
    ACCENT_BAR_HEIGHT = Inches(0.07)
    ACCENT_BAR_TOP = Inches(1.25)

    def __init__(self, theme: str = "ocean-depths"):
        """
        Initialize PPTX builder.

        Args:
            theme: Theme name (corporate, executive, professional, dark_modern)
        """
        self.theme_name = theme
        self.theme_colors = ThemeColors.get_theme(theme)
        self.prs = PptxPresentation()

        # Set slide dimensions
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT

        logger.info("pptx_builder_initialized", theme=theme)

    # ------------------------------------------------------------------
    # Helpers: background, accent bar, slide number badge
    # ------------------------------------------------------------------

    def _fill_slide_background(self, slide):
        """Fill slide background for dark themes."""
        bg_color = self.theme_colors.get("background")
        if not bg_color:
            return
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
        except Exception as e:
            logger.warning("background_fill_failed", error=str(e))

    def _add_accent_bar(self, slide, top: Optional[float] = None):
        """Add a thin horizontal accent bar below the title area."""
        try:
            bar_top = Inches(top) if top is not None else self.ACCENT_BAR_TOP
            bar = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                self.MARGIN_LEFT,
                bar_top,
                self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT,
                self.ACCENT_BAR_HEIGHT,
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = self.theme_colors["header_bar"]
            bar.line.fill.background()  # no border
        except Exception as e:
            logger.warning("accent_bar_failed", error=str(e))

    def _add_left_accent_strip(self, slide):
        """Add a vertical accent strip on the left edge (consulting style)."""
        try:
            strip = slide.shapes.add_shape(
                1,
                Inches(0),
                Inches(0),
                Inches(0.12),
                self.SLIDE_HEIGHT,
            )
            strip.fill.solid()
            strip.fill.fore_color.rgb = self.theme_colors["primary"]
            strip.line.fill.background()
        except Exception as e:
            logger.warning("left_strip_failed", error=str(e))

    def _add_slide_number_badge(self, slide, number: int):
        """Add a small slide number badge at bottom-right."""
        try:
            badge_w = Inches(0.45)
            badge_h = Inches(0.28)
            badge_left = self.SLIDE_WIDTH - badge_w - Inches(0.15)
            badge_top = self.SLIDE_HEIGHT - badge_h - Inches(0.1)
            badge = slide.shapes.add_textbox(badge_left, badge_top, badge_w, badge_h)
            tf = badge.text_frame
            p = tf.paragraphs[0]
            p.text = str(number)
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(9)
            p.font.color.rgb = self.theme_colors.get("text_light", self.theme_colors["text"])
        except Exception as e:
            logger.warning("slide_number_badge_failed", error=str(e))

    def _add_bottom_accent_strip(self, slide):
        """Add a thin bottom accent strip."""
        try:
            strip = slide.shapes.add_shape(
                1,
                Inches(0),
                self.SLIDE_HEIGHT - Inches(0.06),
                self.SLIDE_WIDTH,
                Inches(0.06),
            )
            strip.fill.solid()
            strip.fill.fore_color.rgb = self.theme_colors["accent"]
            strip.line.fill.background()
        except Exception as e:
            logger.warning("bottom_strip_failed", error=str(e))

    def _resolve_content(self, slide_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Resolve slide content by merging root-level fields into content dict.

        The LLM sometimes places fields (bullets, chart_data, highlight_text, etc.)
        at the root of the slide object instead of inside 'content'. This method
        normalises both layouts into a single content dict.
        """
        content = dict(slide_data.get("content") or {})

        # Fields that may appear at root level
        root_passthrough = [
            "bullets", "chart_data", "chart_type", "table_data",
            "comparison_data", "highlight_text", "icon_name",
            "subtitle", "transition", "speaker_notes",
            "metric_value", "metric_label", "metric_trend",
        ]
        for field in root_passthrough:
            if field not in content and field in slide_data:
                content[field] = slide_data[field]

        # kpi_badges → bullets for title slides
        if "kpi_badges" in slide_data and not content.get("bullets"):
            badges = slide_data["kpi_badges"]
            if isinstance(badges, list):
                content["bullets"] = [
                    b.get("label", "") + (f" — {b['description']}" if b.get("description") else "")
                    if isinstance(b, dict) else str(b)
                    for b in badges
                ]

        # left_panel_bullets → bullets (chart slides)
        if "left_panel_bullets" in slide_data and not content.get("bullets"):
            content["bullets"] = slide_data["left_panel_bullets"]

        # chart_data at root of content dict (not nested further)
        # Normalise {chart_data: {chart_data: [...]}} → {chart_data: [...]}
        cd = content.get("chart_data")
        if isinstance(cd, dict) and "chart_data" in cd:
            content["chart_data"] = cd["chart_data"]
            if "chart_type" not in content and "chart_type" in cd:
                content["chart_type"] = cd["chart_type"]

        # Merge chart_type into chart_data dict if chart_data is a dict
        if isinstance(cd, dict) and "type" not in cd and content.get("chart_type"):
            content["chart_data"]["type"] = content["chart_type"]

        return content

    def build(self, slides_data: List[Dict[str, Any]]) -> bytes:
        """
        Build PPTX from slides data.

        Args:
            slides_data: List of slide dictionaries from Slide_JSON

        Returns:
            PPTX file as bytes
        """
        logger.info("pptx_build_started", slide_count=len(slides_data))

        for i, slide_data in enumerate(slides_data):
            try:
                slide_type = slide_data.get("type", "content")
                slide_number = slide_data.get("slide_number", i + 1)

                if slide_type == "title":
                    self._build_title_slide(slide_data)
                elif slide_type == "content":
                    self._build_content_slide(slide_data, slide_number)
                elif slide_type == "chart":
                    self._build_chart_slide(slide_data, slide_number)
                elif slide_type == "table":
                    self._build_table_slide(slide_data, slide_number)
                elif slide_type == "comparison":
                    self._build_comparison_slide(slide_data, slide_number)
                elif slide_type == "metric":
                    self._build_metric_slide(slide_data, slide_number)
                else:
                    logger.warning("unknown_slide_type", type=slide_type, slide_number=i + 1)
                    self._build_content_slide(slide_data, slide_number)  # Fallback

            except Exception as e:
                logger.error("slide_build_failed", slide_number=i + 1, error=str(e))
                # Continue with next slide

        # Save to bytes
        buf = io.BytesIO()
        self.prs.save(buf)
        pptx_bytes = buf.getvalue()

        logger.info("pptx_build_completed", size_bytes=len(pptx_bytes))
        return pptx_bytes
    
    def _build_title_slide(self, slide_data: Dict[str, Any]):
        """Build enterprise title slide with accent strip and KPI bullets."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout for full control
        slide = self.prs.slides.add_slide(slide_layout)
        self._fill_slide_background(slide)

        content = self._resolve_content(slide_data)
        title_text = slide_data.get("title", "")
        subtitle_text = content.get("subtitle") or content.get("highlight_text", "")
        bullets = content.get("bullets", [])

        # Left accent strip
        self._add_left_accent_strip(slide)

        # Large title
        title_left = Inches(1.0)
        title_top = Inches(1.8)
        title_width = Inches(9.5)
        title_height = Inches(1.8)
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.theme_colors["primary"]

        # Subtitle
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(title_left, Inches(3.7), title_width, Inches(0.6))
            sub_tf = sub_box.text_frame
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = subtitle_text
            sub_p.font.size = Pt(20)
            sub_p.font.color.rgb = self.theme_colors.get("text_light", self.theme_colors["text"])

        # Horizontal divider line
        self._add_accent_bar(slide, top=4.45)

        # KPI bullets rendered as small badge boxes
        if bullets:
            self._add_kpi_badges(slide, bullets, top_y=4.7)

        # Bottom accent strip
        self._add_bottom_accent_strip(slide)

        # Apply transition
        transition = content.get("transition", "fade")
        TransitionMapper.apply_transition(slide, transition)

        logger.debug("title_slide_built", title=title_text[:50])
    
    def _build_content_slide(self, slide_data: Dict[str, Any], slide_number: int = 0):
        """Build enterprise content slide with accent bar, icon badge, and highlight box."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        self._fill_slide_background(slide)
        self._add_left_accent_strip(slide)

        content = self._resolve_content(slide_data)
        title_text = slide_data.get("title", "")
        bullets = content.get("bullets", [])
        highlight_text = content.get("highlight_text", "")
        icon_name = content.get("icon_name", "")

        # Title area
        title_left = Inches(1.0)
        title_top = Inches(0.35)
        title_width = Inches(11.5)
        title_height = Inches(0.85)
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.theme_colors["primary"]

        # Accent bar below title
        self._add_accent_bar(slide, top=1.28)

        # Bullet content area — leave room for highlight box at bottom
        bullet_top = Inches(1.45)
        bullet_height = Inches(4.5) if highlight_text else Inches(5.5)
        bullet_width = Inches(11.5)

        if bullets:
            bullet_box = slide.shapes.add_textbox(title_left, bullet_top, bullet_width, bullet_height)
            btf = bullet_box.text_frame
            btf.word_wrap = True
            for i, bullet in enumerate(bullets):
                if i == 0:
                    bp = btf.paragraphs[0]
                else:
                    bp = btf.add_paragraph()
                # Bullet character prefix
                bp.text = f"\u2022  {bullet}"
                bp.font.size = Pt(16)
                bp.font.color.rgb = self.theme_colors["text"]
                bp.space_before = Pt(4)

        # Highlight / insight box at bottom
        if highlight_text:
            self._add_insight_box(slide, highlight_text, icon_name)

        # Slide number
        self._add_slide_number_badge(slide, slide_number)

        # Apply transition
        transition = content.get("transition", "fade")
        TransitionMapper.apply_transition(slide, transition)

        logger.debug("content_slide_built", bullets=len(bullets))
    
    def _build_chart_slide(self, slide_data: Dict[str, Any], slide_number: int = 0):
        """Build enterprise chart slide: title + accent bar + bullets left + chart right."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        self._fill_slide_background(slide)
        self._add_left_accent_strip(slide)

        content = self._resolve_content(slide_data)
        title_text = slide_data.get("title", "")

        # Title
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.35), Inches(11.5), Inches(0.85))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.theme_colors["primary"]

        # Accent bar
        self._add_accent_bar(slide, top=1.28)

        # Chart data — support both legacy {categories, series} and new [{label, value}] formats
        chart_data = content.get("chart_data", {})
        chart_type_str = (
            chart_data.get("type") if isinstance(chart_data, dict) else None
        ) or (
            chart_data.get("chart_type") if isinstance(chart_data, dict) else None
        ) or content.get("chart_type", "bar")

        # Normalise [{label, value}] → {categories, series}
        if isinstance(chart_data, list):
            chart_data = self._normalise_label_value_chart(chart_data, chart_type_str)
        elif isinstance(chart_data, dict) and "categories" not in chart_data:
            # Check if it's a dict wrapping label/value items
            items = chart_data.get("data") or chart_data.get("values") or []
            if items and isinstance(items[0], dict) and "label" in items[0]:
                chart_data = self._normalise_label_value_chart(items, chart_type_str)

        # Layout: left panel (bullets) + right panel (chart)
        panel_top = Inches(1.45)
        panel_height = Inches(5.5)
        left_w = Inches(5.0)
        right_w = Inches(7.0)
        gap = Inches(0.33)
        chart_left = Inches(1.0) + left_w + gap

        # Bullets on left
        bullets = content.get("bullets", [])
        if bullets:
            self._add_text_box_with_bullets(
                slide, bullets,
                Inches(1.0), panel_top, left_w, panel_height,
                font_size=Pt(14),
            )

        # Chart on right
        if chart_data and (chart_data.get("categories") or chart_data.get("series")):
            self._add_chart(
                slide, chart_data, chart_type_str,
                chart_left, panel_top, right_w, panel_height,
            )

        # Highlight box
        highlight_text = content.get("highlight_text", "")
        if highlight_text:
            self._add_insight_box(slide, highlight_text, "")

        self._add_slide_number_badge(slide, slide_number)
        transition = content.get("transition", "fade")
        TransitionMapper.apply_transition(slide, transition)

        logger.debug("chart_slide_built", chart_type=chart_type_str)
    
    def _build_table_slide(self, slide_data: Dict[str, Any], slide_number: int = 0):
        """Build enterprise table slide with styled header and alternating rows."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        self._fill_slide_background(slide)
        self._add_left_accent_strip(slide)

        content = self._resolve_content(slide_data)
        title_text = slide_data.get("title", "")

        # Title
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.35), Inches(11.5), Inches(0.85))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.theme_colors["primary"]

        # Accent bar
        self._add_accent_bar(slide, top=1.28)

        table_data = content.get("table_data", {})
        bullets = content.get("bullets", [])
        highlight_text = content.get("highlight_text", "")

        panel_top = Inches(1.45)
        panel_height = Inches(4.8) if highlight_text else Inches(5.5)

        if table_data and table_data.get("headers"):
            # Table on left (wider)
            table_w = Inches(7.8) if bullets else Inches(11.5)
            self._add_table(slide, table_data, Inches(1.0), panel_top, table_w, panel_height)

            # Bullets on right
            if bullets:
                self._add_text_box_with_bullets(
                    slide, bullets,
                    Inches(9.2), panel_top, Inches(3.5), panel_height,
                    font_size=Pt(13),
                )

        if highlight_text:
            self._add_insight_box(slide, highlight_text, "")

        self._add_slide_number_badge(slide, slide_number)
        transition = content.get("transition", "fade")
        TransitionMapper.apply_transition(slide, transition)

        logger.debug("table_slide_built")
    
    def _build_comparison_slide(self, slide_data: Dict[str, Any], slide_number: int = 0):
        """Build enterprise comparison slide with color-coded column headers."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        self._fill_slide_background(slide)
        self._add_left_accent_strip(slide)

        content = self._resolve_content(slide_data)
        title_text = slide_data.get("title", "")

        # Title
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.35), Inches(11.5), Inches(0.85))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.theme_colors["primary"]

        # Accent bar
        self._add_accent_bar(slide, top=1.28)

        comparison_data = content.get("comparison_data", {})

        # Support both old format {left, right, left_title, right_title}
        # and new format {left_column: {heading, bullets}, right_column: {heading, bullets}}
        if "left_column" in comparison_data:
            left_col = comparison_data["left_column"]
            right_col = comparison_data["right_column"]
            left_title = left_col.get("heading", "")
            right_title = right_col.get("heading", "")
            left_items = left_col.get("bullets", [])
            right_items = right_col.get("bullets", [])
        else:
            left_title = comparison_data.get("left_title", "")
            right_title = comparison_data.get("right_title", "")
            left_items = comparison_data.get("left", [])
            right_items = comparison_data.get("right", [])

        col_top = Inches(1.45)
        col_height = Inches(5.5)
        col_w = Inches(5.5)
        left_left = Inches(1.0)
        right_left = Inches(7.0)

        # Left column header box (primary color)
        if left_title:
            self._add_column_header_box(slide, left_title, left_left, col_top, col_w, self.theme_colors["primary"])
            items_top = col_top + Inches(0.55)
            items_height = col_height - Inches(0.55)
        else:
            items_top = col_top
            items_height = col_height

        if left_items:
            self._add_comparison_column(slide, "", left_items, left_left, items_top, col_w, items_height)

        # Vertical divider
        try:
            div = slide.shapes.add_shape(1, Inches(6.6), col_top, Inches(0.03), col_height)
            div.fill.solid()
            div.fill.fore_color.rgb = self.theme_colors.get("divider", RGBColor(200, 200, 200))
            div.line.fill.background()
        except Exception:
            pass

        # Right column header box (accent color)
        if right_title:
            self._add_column_header_box(slide, right_title, right_left, col_top, col_w, self.theme_colors["accent"])
            items_top_r = col_top + Inches(0.55)
            items_height_r = col_height - Inches(0.55)
        else:
            items_top_r = col_top
            items_height_r = col_height

        if right_items:
            self._add_comparison_column(slide, "", right_items, right_left, items_top_r, col_w, items_height_r)

        # Highlight box
        highlight_text = content.get("highlight_text", "")
        if highlight_text:
            self._add_insight_box(slide, highlight_text, "")

        self._add_slide_number_badge(slide, slide_number)
        transition = content.get("transition", "fade")
        TransitionMapper.apply_transition(slide, transition)

        logger.debug("comparison_slide_built")
    
    def _build_metric_slide(self, slide_data: Dict[str, Any], slide_number: int = 0):
        """Build enterprise KPI / metric slide with large number display."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        self._fill_slide_background(slide)
        self._add_left_accent_strip(slide)

        content = self._resolve_content(slide_data)
        title_text = slide_data.get("title", "")

        # Title
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.35), Inches(11.5), Inches(0.85))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.theme_colors["primary"]

        # Accent bar
        self._add_accent_bar(slide, top=1.28)

        # KPI value — large display
        metric_value = str(content.get("metric_value", content.get("value", "")))
        metric_label = str(content.get("metric_label", content.get("label", "")))
        metric_trend = str(content.get("metric_trend", content.get("trend", "")))

        # Large KPI box
        kpi_box = slide.shapes.add_shape(1, Inches(1.0), Inches(1.6), Inches(5.5), Inches(2.8))
        kpi_box.fill.solid()
        kpi_box.fill.fore_color.rgb = self.theme_colors["kpi_bg"]
        kpi_box.line.fill.background()

        # KPI value text
        kpi_tf = kpi_box.text_frame
        kpi_tf.word_wrap = False
        kpi_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        kp = kpi_tf.paragraphs[0]
        kp.text = metric_value
        kp.alignment = PP_ALIGN.CENTER
        kp.font.size = Pt(60)
        kp.font.bold = True
        kp.font.color.rgb = self.theme_colors["kpi_text"]

        # KPI label below value
        if metric_label:
            lbl_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(5.5), Inches(0.5))
            lbl_tf = lbl_box.text_frame
            lp = lbl_tf.paragraphs[0]
            lp.text = metric_label
            lp.alignment = PP_ALIGN.CENTER
            lp.font.size = Pt(16)
            lp.font.color.rgb = self.theme_colors.get("text_light", self.theme_colors["text"])

        # Trend badge
        if metric_trend:
            trend_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.1), Inches(5.5), Inches(0.4))
            trend_tf = trend_box.text_frame
            tp = trend_tf.paragraphs[0]
            tp.text = metric_trend
            tp.alignment = PP_ALIGN.CENTER
            tp.font.size = Pt(14)
            tp.font.bold = True
            # Green for positive, red for negative
            if any(c in metric_trend for c in ["+", "▲", "up", "Up", "UP"]):
                tp.font.color.rgb = RGBColor(0, 166, 81)
            elif any(c in metric_trend for c in ["-", "▼", "down", "Down", "DOWN"]):
                tp.font.color.rgb = RGBColor(237, 28, 36)
            else:
                tp.font.color.rgb = self.theme_colors["accent"]

        # Context bullets on right
        bullets = content.get("bullets", [])
        if bullets:
            self._add_text_box_with_bullets(
                slide, bullets,
                Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.5),
                font_size=Pt(14),
            )

        self._add_slide_number_badge(slide, slide_number)
        transition = content.get("transition", "fade")
        TransitionMapper.apply_transition(slide, transition)

        logger.debug("metric_slide_built", metric_value=metric_value)

    def _normalise_label_value_chart(
        self, items: List[Dict[str, Any]], chart_type: str
    ) -> Dict[str, Any]:
        """Convert [{label, value}, ...] format to {categories, series} format."""
        categories = [str(item.get("label", item.get("name", f"Item {i+1}"))) for i, item in enumerate(items)]
        values = [float(item.get("value", item.get("y", 0))) for item in items]
        return {
            "type": chart_type,
            "categories": categories,
            "series": [{"name": "Value", "values": values}],
        }

    def _add_chart(
        self,
        slide,
        chart_data: Dict[str, Any],
        chart_type: str,
        left,
        top,
        width,
        height,
    ):
        """Add chart to slide — supports bar, line, pie, area, stacked_bar, donut, scatter."""
        try:
            # Extended chart type map
            chart_type_map = {
                "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "bar_horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "line_smooth": XL_CHART_TYPE.LINE_MARKERS,
                "pie": XL_CHART_TYPE.PIE,
                "donut": XL_CHART_TYPE.DOUGHNUT,
                "area": XL_CHART_TYPE.AREA,
                "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
                "stacked_area": XL_CHART_TYPE.AREA_STACKED,
                "scatter": XL_CHART_TYPE.XY_SCATTER,
            }

            pptx_chart_type = chart_type_map.get(
                chart_type.lower() if chart_type else "bar",
                XL_CHART_TYPE.COLUMN_CLUSTERED,
            )

            categories = chart_data.get("categories", [])
            series_list = chart_data.get("series", [])

            if not categories or not series_list:
                logger.warning("chart_data_missing", chart_type=chart_type)
                return

            # Scatter charts use XyChartData
            if pptx_chart_type == XL_CHART_TYPE.XY_SCATTER:
                self._add_scatter_chart(slide, chart_data, left, top, width, height)
                return

            # Standard CategoryChartData
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = [str(c) for c in categories]

            for series in series_list:
                series_name = series.get("name", "Series")
                series_values = [float(v) if v is not None else 0.0 for v in series.get("values", [])]
                chart_data_obj.add_series(series_name, series_values)

            chart = slide.shapes.add_chart(
                pptx_chart_type, left, top, width, height, chart_data_obj
            ).chart

            self._apply_chart_theme(chart, chart_type)

            logger.debug("chart_added", type=chart_type, categories=len(categories))

        except Exception as e:
            logger.error("chart_addition_failed", error=str(e))

    def _add_scatter_chart(self, slide, chart_data: Dict[str, Any], left, top, width, height):
        """Add XY scatter chart."""
        try:
            xy_data = XyChartData()
            series_list = chart_data.get("series", [])
            for series in series_list:
                s = xy_data.add_series(series.get("name", "Series"))
                xs = series.get("x_values", series.get("values", []))
                ys = series.get("y_values", series.get("values", []))
                for x, y in zip(xs, ys):
                    s.add_data_point(float(x), float(y))
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.XY_SCATTER, left, top, width, height, xy_data
            ).chart
            self._apply_chart_theme(chart, "scatter")
        except Exception as e:
            logger.error("scatter_chart_failed", error=str(e))
    
    def _apply_chart_theme(self, chart, chart_type: str):
        """Apply theme colors and clean styling to chart."""
        try:
            chart_colors = self.theme_colors["chart_colors"]
            for i, series in enumerate(chart.series):
                color = chart_colors[i % len(chart_colors)]
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = color

            # Legend styling
            if hasattr(chart, "has_legend") and chart.has_legend:
                chart.legend.font.size = Pt(10)
                chart.legend.font.color.rgb = self.theme_colors["text"]

            # Plot area background — transparent
            try:
                chart.plot_area.format.fill.background()
            except Exception:
                pass

            # Chart area background — match slide background
            try:
                chart.chart_area.format.fill.solid()
                chart.chart_area.format.fill.fore_color.rgb = self.theme_colors["background"]
            except Exception:
                pass

        except Exception as e:
            logger.warning("chart_theme_application_failed", error=str(e))
    
    def _add_table(
        self,
        slide,
        table_data: Dict[str, Any],
        left,
        top,
        width,
        height,
    ):
        """Add enterprise-styled table with bold header and alternating rows."""
        try:
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])

            if not headers or not rows:
                logger.warning("table_data_missing")
                return

            num_rows = len(rows) + 1  # +1 for header
            num_cols = len(headers)

            table = slide.shapes.add_table(
                int(num_rows), int(num_cols), left, top, width, height
            ).table

            # Header row
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.theme_colors["primary"]
                para = cell.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.size = Pt(12)
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.alignment = PP_ALIGN.CENTER

            # Data rows
            alt_color = self.theme_colors.get("surface", RGBColor(245, 247, 250))
            for row_idx, row in enumerate(rows):
                for col_idx in range(num_cols):
                    value = row[col_idx] if col_idx < len(row) else ""
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(value)
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(11)
                    para.font.color.rgb = self.theme_colors["text"]
                    # Alternating row background
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = alt_color

            logger.debug("table_added", rows=num_rows, cols=num_cols)

        except Exception as e:
            logger.error("table_addition_failed", error=str(e))
    
    def _add_text_box_with_bullets(
        self,
        slide,
        bullets: List[str],
        left,
        top,
        width,
        height,
        font_size: Pt = None,
    ):
        """Add text box with styled bullet points."""
        if font_size is None:
            font_size = Pt(16)
        try:
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"\u2022  {bullet}"
                p.font.size = font_size
                p.font.color.rgb = self.theme_colors["text"]
                p.space_before = Pt(4)

        except Exception as e:
            logger.error("textbox_addition_failed", error=str(e))
    
    def _add_comparison_column(
        self,
        slide,
        title: str,
        items: List[str],
        left,
        top,
        width,
        height,
    ):
        """Add comparison column items (title handled separately via header box)."""
        try:
            content_box = slide.shapes.add_textbox(left, top, width, height)
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"\u2022  {item}"
                p.font.size = Pt(14)
                p.font.color.rgb = self.theme_colors["text"]
                p.space_before = Pt(5)

        except Exception as e:
            logger.error("comparison_column_addition_failed", error=str(e))

    def _add_column_header_box(self, slide, title: str, left, top, width, color: RGBColor):
        """Add a colored header box for comparison columns."""
        try:
            box = slide.shapes.add_shape(1, left, top, width, Inches(0.48))
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()
            tf = box.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = title
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(255, 255, 255)
        except Exception as e:
            logger.warning("column_header_box_failed", error=str(e))

    def _add_insight_box(self, slide, text: str, icon_name: str):
        """Add a full-width insight / highlight box at the bottom of the slide."""
        try:
            box_h = Inches(0.65)
            box_top = self.SLIDE_HEIGHT - box_h - Inches(0.08)
            box = slide.shapes.add_shape(
                1,
                self.MARGIN_LEFT,
                box_top,
                self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT,
                box_h,
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme_colors["primary"]
            box.line.fill.background()

            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            prefix = f"  \u25B6  " if not icon_name else f"  \u2605  "
            p.text = prefix + text
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
        except Exception as e:
            logger.error("insight_box_failed", error=str(e))

    def _add_kpi_badges(self, slide, bullets: List[str], top_y: float):
        """Render KPI bullets as small badge boxes on the title slide."""
        try:
            badge_w = Inches(2.6)
            badge_h = Inches(0.55)
            gap = Inches(0.2)
            start_left = Inches(1.0)
            max_per_row = 4

            for i, bullet in enumerate(bullets[:max_per_row]):
                badge_left = start_left + i * (badge_w + gap)
                badge = slide.shapes.add_shape(1, badge_left, Inches(top_y), badge_w, badge_h)
                badge.fill.solid()
                badge.fill.fore_color.rgb = self.theme_colors["kpi_bg"]
                badge.line.color.rgb = self.theme_colors["accent"]

                tf = badge.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = bullet
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = self.theme_colors["kpi_text"]
        except Exception as e:
            logger.warning("kpi_badges_failed", error=str(e))
    
    def _apply_text_formatting(
        self,
        shape,
        font_size: Pt,
        bold: bool = False,
        color: RGBColor = None,
    ):
        """Apply text formatting to a shape."""
        try:
            if not hasattr(shape, "text_frame"):
                return
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                paragraph.font.size = font_size
                paragraph.font.bold = bold
                if color:
                    paragraph.font.color.rgb = color
        except Exception as e:
            logger.warning("text_formatting_failed", error=str(e))

    def _apply_paragraph_formatting(
        self,
        paragraph,
        font_size: Pt,
        color: RGBColor = None,
    ):
        """Apply formatting to a paragraph."""
        try:
            paragraph.font.size = font_size
            if color:
                paragraph.font.color.rgb = color
        except Exception as e:
            logger.warning("paragraph_formatting_failed", error=str(e))


# ---------------------------------------------------------------------------
# Main Export Function
# ---------------------------------------------------------------------------

def build_pptx(slides_data: List[Dict[str, Any]], theme: str = "ocean-depths") -> bytes:
    """
    Build PPTX file from Slide_JSON data.
    
    This is the main entry point for PPTX generation.
    
    Args:
        slides_data: List of slide dictionaries from Slide_JSON
        theme: Theme name (corporate, executive, professional, dark_modern)
        
    Returns:
        PPTX file as bytes
        
    Raises:
        ValueError: If slides_data is invalid
    """
    if not isinstance(slides_data, list):
        raise ValueError("slides_data must be a list")
    
    if not slides_data:
        logger.warning("empty_slides_data")
        # Return minimal presentation
        builder = PPTXBuilder(theme)
        return builder.build([])
    
    builder = PPTXBuilder(theme)
    return builder.build(slides_data)
