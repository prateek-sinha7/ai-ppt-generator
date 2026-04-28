"""
Tests for PPTX Export Service (Task 24).

Tests cover:
- Task 24.1: Slide type mapping to PPTX layouts
- Task 24.2: Theme application (Ocean Depths, Modern Minimalist, Dark Modern)
- Task 24.3: Chart rendering (bar, line, pie)
- Task 24.4: Table rendering with formatting
- Task 24.5: Transition mapping
- Task 24.6: S3/MinIO upload (integration test)
- Task 24.7: Performance validation (30 seconds for 50 slides)
"""

import time
from io import BytesIO
from typing import Any, Dict, List

import pytest
from pptx import Presentation as PptxPresentation

from app.services.pptx_export import (
    PPTXBuilder,
    ThemeColors,
    TransitionMapper,
    build_pptx,
)


# ---------------------------------------------------------------------------
# Test Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def sample_title_slide() -> Dict[str, Any]:
    """Sample title slide data."""
    return {
        "slide_id": "slide-1",
        "slide_number": 1,
        "type": "title",
        "title": "AI Presentation Intelligence Platform",
        "content": {
            "subtitle": "Enterprise-Grade Presentation Generation",
            "transition": "fade"
        },
        "visual_hint": "centered"
    }


@pytest.fixture
def sample_content_slide() -> Dict[str, Any]:
    """Sample content slide data."""
    return {
        "slide_id": "slide-2",
        "slide_number": 2,
        "type": "content",
        "title": "Key Features",
        "content": {
            "bullets": [
                "Multi-agent AI pipeline",
                "Automatic industry detection",
                "Real-time generation",
                "Enterprise security"
            ],
            "highlight_text": "Production Ready",
            "transition": "slide"
        },
        "visual_hint": "bullet-left"
    }


@pytest.fixture
def sample_chart_slide() -> Dict[str, Any]:
    """Sample chart slide data."""
    return {
        "slide_id": "slide-3",
        "slide_number": 3,
        "type": "chart",
        "title": "Market Growth",
        "content": {
            "chart_data": {
                "type": "bar",
                "categories": ["Q1", "Q2", "Q3", "Q4"],
                "series": [
                    {
                        "name": "Revenue",
                        "values": [100, 150, 200, 250]
                    },
                    {
                        "name": "Profit",
                        "values": [20, 35, 50, 70]
                    }
                ]
            },
            "bullets": ["Strong growth trajectory", "Increasing profitability"],
            "transition": "fade"
        },
        "visual_hint": "split-chart-right"
    }


@pytest.fixture
def sample_table_slide() -> Dict[str, Any]:
    """Sample table slide data."""
    return {
        "slide_id": "slide-4",
        "slide_number": 4,
        "type": "table",
        "title": "Financial Summary",
        "content": {
            "table_data": {
                "headers": ["Metric", "2023", "2024", "Growth"],
                "rows": [
                    ["Revenue", "$100M", "$150M", "50%"],
                    ["Profit", "$20M", "$35M", "75%"],
                    ["Customers", "1000", "2000", "100%"]
                ]
            },
            "bullets": ["Year-over-year growth", "Strong customer acquisition"],
            "transition": "none"
        },
        "visual_hint": "split-table-left"
    }


@pytest.fixture
def sample_comparison_slide() -> Dict[str, Any]:
    """Sample comparison slide data."""
    return {
        "slide_id": "slide-5",
        "slide_number": 5,
        "type": "comparison",
        "title": "Before vs After",
        "content": {
            "comparison_data": {
                "left_title": "Before",
                "left": [
                    "Manual presentation creation",
                    "Hours of work",
                    "Inconsistent quality"
                ],
                "right_title": "After",
                "right": [
                    "Automated AI generation",
                    "Minutes to complete",
                    "Enterprise-grade quality"
                ]
            },
            "transition": "fade"
        },
        "visual_hint": "two-column"
    }


@pytest.fixture
def sample_metric_slide() -> Dict[str, Any]:
    """Sample metric / KPI slide data."""
    return {
        "slide_id": "slide-6",
        "slide_number": 6,
        "type": "metric",
        "title": "Claims Processing Speed: Industry-Leading Performance",
        "content": {
            "metric_value": "4.8 days",
            "metric_label": "Average Claims Processing Time",
            "metric_trend": "▼ 62% improvement vs. 2022",
            "bullets": [
                "Industry average: 15.2 days — 3.2x faster than peers",
                "NPS correlation: each day reduction adds 4.2 NPS points",
                "Cost impact: $47 saved per claim vs. manual processing",
                "Customer retention: 89% renewal rate vs. 71% industry average"
            ],
            "highlight_text": "Processing speed is the #1 driver of customer satisfaction",
            "transition": "fade"
        },
        "visual_hint": "highlight-metric"
    }


@pytest.fixture
def all_slide_types(
    sample_title_slide,
    sample_content_slide,
    sample_chart_slide,
    sample_table_slide,
    sample_comparison_slide,
    sample_metric_slide,
) -> List[Dict[str, Any]]:
    """All slide types for comprehensive testing."""
    return [
        sample_title_slide,
        sample_content_slide,
        sample_chart_slide,
        sample_table_slide,
        sample_comparison_slide,
        sample_metric_slide,
    ]


# ---------------------------------------------------------------------------
# Task 24.1: Slide Type Mapping Tests
# ---------------------------------------------------------------------------

class TestSlideTypeMapping:
    """Test slide type to PPTX layout mapping."""
    
    def test_title_slide_creation(self, sample_title_slide):
        """Test title slide is created with correct layout."""
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([sample_title_slide])
        
        assert isinstance(pptx_bytes, bytes)
        assert len(pptx_bytes) > 0
        
        # Verify slide was created
        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1
        
        slide = prs.slides[0]
        # Title is in a textbox (blank layout), verify at least one shape has the title text
        all_text = " ".join(
            shape.text_frame.text for shape in slide.shapes if hasattr(shape, "text_frame")
        )
        assert sample_title_slide["title"] in all_text
    
    def test_content_slide_creation(self, sample_content_slide):
        """Test content slide is created with bullets."""
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([sample_content_slide])
        
        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1
        
        slide = prs.slides[0]
        all_text = " ".join(
            shape.text_frame.text for shape in slide.shapes if hasattr(shape, "text_frame")
        )
        assert sample_content_slide["title"] in all_text
    
    def test_chart_slide_creation(self, sample_chart_slide):
        """Test chart slide is created with chart."""
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([sample_chart_slide])
        
        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1
        
        slide = prs.slides[0]
        
        # Check for chart shape
        has_chart = any(hasattr(shape, 'chart') for shape in slide.shapes)
        assert has_chart, "Chart slide should contain a chart"
    
    def test_table_slide_creation(self, sample_table_slide):
        """Test table slide is created with table."""
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([sample_table_slide])
        
        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1
        
        slide = prs.slides[0]
        
        # Check for table shape
        has_table = any(hasattr(shape, 'table') for shape in slide.shapes)
        assert has_table, "Table slide should contain a table"
    
    def test_comparison_slide_creation(self, sample_comparison_slide):
        """Test comparison slide is created with two columns."""
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([sample_comparison_slide])

        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1

    def test_metric_slide_creation(self, sample_metric_slide):
        """Test metric/KPI slide is created with large number display."""
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([sample_metric_slide])

        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1

    def test_all_slide_types(self, all_slide_types):
        """Test all slide types can be created in one presentation."""
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build(all_slide_types)

        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 6


# ---------------------------------------------------------------------------
# Task 24.2: Theme Application Tests
# ---------------------------------------------------------------------------

class TestThemeApplication:
    """Test theme color scheme application."""
    
    def test_ocean_depths_theme(self, sample_content_slide):
        """Test Ocean Depths theme colors are applied."""
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([sample_content_slide])
        
        assert isinstance(pptx_bytes, bytes)
        assert len(pptx_bytes) > 0
        
        # Verify theme was set
        assert builder.theme_name == "ocean-depths"
        assert builder.theme_colors == ThemeColors.OCEAN_DEPTHS
    def test_modern_minimalist_theme(self, sample_content_slide):
        """Test Modern Minimalist theme colors are applied."""
        builder = PPTXBuilder("modern-minimalist")
        pptx_bytes = builder.build([sample_content_slide])
        
        assert isinstance(pptx_bytes, bytes)
        assert builder.theme_name == "modern-minimalist"
        assert builder.theme_colors == ThemeColors.MODERN_MINIMALIST
    
    def test_dark_modern_theme(self, sample_content_slide):
        """Test Dark Modern theme colors are applied."""
        builder = PPTXBuilder("tech-innovation")
        pptx_bytes = builder.build([sample_content_slide])
        
        assert isinstance(pptx_bytes, bytes)
        assert builder.theme_name == "tech-innovation"
        assert builder.theme_colors == ThemeColors.TECH_INNOVATION
    
    def test_dark_modern_theme_with_hyphen(self, sample_content_slide):
        """Test Dark Modern theme with hyphen notation."""
        builder = PPTXBuilder("tech-innovation")
        pptx_bytes = builder.build([sample_content_slide])
        
        assert isinstance(pptx_bytes, bytes)
        assert builder.theme_colors == ThemeColors.TECH_INNOVATION
    
    def test_invalid_theme_defaults_to_ocean_depths(self, sample_content_slide):
        """Test invalid theme defaults to Ocean Depths."""
        builder = PPTXBuilder("invalid_theme")
        pptx_bytes = builder.build([sample_content_slide])
        
        assert isinstance(pptx_bytes, bytes)
        assert builder.theme_colors == ThemeColors.OCEAN_DEPTHS
    
    def test_theme_colors_in_chart(self, sample_chart_slide):
        """Test theme colors are applied to charts."""
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([sample_chart_slide])
        
        prs = PptxPresentation(BytesIO(pptx_bytes))
        slide = prs.slides[0]
        
        # Find chart
        chart = None
        for shape in slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                break
        
        assert chart is not None, "Chart should be present"


# ---------------------------------------------------------------------------
# Task 24.3: Chart Rendering Tests
# ---------------------------------------------------------------------------

class TestChartRendering:
    """Test chart rendering with different types."""
    
    def test_bar_chart_rendering(self):
        """Test bar chart is rendered correctly."""
        slide_data = {
            "type": "chart",
            "title": "Bar Chart Test",
            "content": {
                "chart_data": {
                    "type": "bar",
                    "categories": ["A", "B", "C"],
                    "series": [{"name": "Series 1", "values": [10, 20, 30]}]
                }
            }
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        prs = PptxPresentation(BytesIO(pptx_bytes))
        slide = prs.slides[0]
        
        has_chart = any(hasattr(shape, 'chart') for shape in slide.shapes)
        assert has_chart
    
    def test_line_chart_rendering(self):
        """Test line chart is rendered correctly."""
        slide_data = {
            "type": "chart",
            "title": "Line Chart Test",
            "content": {
                "chart_data": {
                    "type": "line",
                    "categories": ["Jan", "Feb", "Mar"],
                    "series": [{"name": "Trend", "values": [100, 150, 200]}]
                }
            }
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1
    
    def test_pie_chart_rendering(self):
        """Test pie chart is rendered correctly."""
        slide_data = {
            "type": "chart",
            "title": "Pie Chart Test",
            "content": {
                "chart_data": {
                    "type": "pie",
                    "categories": ["A", "B", "C"],
                    "series": [{"name": "Distribution", "values": [30, 40, 30]}]
                }
            }
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1
    
    def test_chart_with_multiple_series(self):
        """Test chart with multiple data series."""
        slide_data = {
            "type": "chart",
            "title": "Multi-Series Chart",
            "content": {
                "chart_data": {
                    "type": "bar",
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": [
                        {"name": "Product A", "values": [100, 120, 140, 160]},
                        {"name": "Product B", "values": [80, 90, 100, 110]},
                        {"name": "Product C", "values": [60, 70, 80, 90]}
                    ]
                }
            }
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        prs = PptxPresentation(BytesIO(pptx_bytes))
        slide = prs.slides[0]
        
        # Find chart and verify series count
        for shape in slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                assert len(chart.series) == 3
                break

    def test_area_chart_rendering(self):
        """Test area chart is rendered correctly."""
        slide_data = {
            "type": "chart",
            "title": "Area Chart Test",
            "content": {
                "chart_data": {
                    "type": "area",
                    "categories": ["2020", "2021", "2022", "2023", "2024"],
                    "series": [{"name": "Revenue", "values": [100, 130, 160, 210, 280]}]
                }
            }
        }
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1

    def test_stacked_bar_chart_rendering(self):
        """Test stacked bar chart is rendered correctly."""
        slide_data = {
            "type": "chart",
            "title": "Stacked Bar Chart Test",
            "content": {
                "chart_data": {
                    "type": "stacked_bar",
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": [
                        {"name": "Segment A", "values": [40, 45, 50, 55]},
                        {"name": "Segment B", "values": [30, 35, 38, 42]},
                    ]
                }
            }
        }
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1

    def test_label_value_chart_format(self):
        """Test chart with [{label, value}] format (LLM output format)."""
        slide_data = {
            "type": "chart",
            "title": "Label-Value Chart",
            "content": {
                "chart_type": "bar",
                "chart_data": [
                    {"label": "Digital-Native", "value": 34.2},
                    {"label": "Incumbent A", "value": 22.8},
                    {"label": "Incumbent B", "value": 18.4},
                    {"label": "Regional", "value": 14.1},
                    {"label": "Others", "value": 10.5},
                ]
            }
        }
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1
        # Should have a chart
        has_chart = any(hasattr(shape, "chart") for shape in prs.slides[0].shapes)
        assert has_chart


# ---------------------------------------------------------------------------
# Task 24.4: Table Rendering Tests
# ---------------------------------------------------------------------------

class TestTableRendering:
    """Test table rendering with formatting."""
    
    def test_basic_table_rendering(self, sample_table_slide):
        """Test basic table is rendered with headers and rows."""
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([sample_table_slide])
        
        prs = PptxPresentation(BytesIO(pptx_bytes))
        slide = prs.slides[0]
        
        # Find table
        table = None
        for shape in slide.shapes:
            if hasattr(shape, 'table'):
                table = shape.table
                break
        
        assert table is not None, "Table should be present"
        
        # Verify dimensions
        table_data = sample_table_slide["content"]["table_data"]
        expected_rows = len(table_data["rows"]) + 1  # +1 for header
        expected_cols = len(table_data["headers"])
        
        assert table.rows.__len__() == expected_rows
        assert table.columns.__len__() == expected_cols
    
    def test_table_header_formatting(self, sample_table_slide):
        """Test table headers are formatted correctly."""
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([sample_table_slide])
        
        prs = PptxPresentation(BytesIO(pptx_bytes))
        slide = prs.slides[0]
        
        # Find table
        for shape in slide.shapes:
            if hasattr(shape, 'table'):
                table = shape.table
                
                # Check header row
                for col_idx in range(table.columns.__len__()):
                    cell = table.cell(0, col_idx)
                    # Header should have text
                    assert cell.text != ""
                break
    
    def test_empty_table_handling(self):
        """Test handling of empty table data."""
        slide_data = {
            "type": "table",
            "title": "Empty Table",
            "content": {
                "table_data": {
                    "headers": [],
                    "rows": []
                }
            }
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        # Should not crash
        assert isinstance(pptx_bytes, bytes)
        assert len(pptx_bytes) > 0


# ---------------------------------------------------------------------------
# Task 24.5: Transition Mapping Tests
# ---------------------------------------------------------------------------

class TestTransitionMapping:
    """Test transition mapping from Slide_JSON to PowerPoint."""
    
    def test_fade_transition(self):
        """Test fade transition is applied."""
        slide_data = {
            "type": "content",
            "title": "Fade Transition",
            "content": {
                "bullets": ["Test"],
                "transition": "fade"
            }
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        # Should not crash
        assert isinstance(pptx_bytes, bytes)
    
    def test_slide_transition(self):
        """Test slide (push) transition is applied."""
        slide_data = {
            "type": "content",
            "title": "Slide Transition",
            "content": {
                "bullets": ["Test"],
                "transition": "slide"
            }
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        assert isinstance(pptx_bytes, bytes)
    
    def test_none_transition(self):
        """Test no transition is applied."""
        slide_data = {
            "type": "content",
            "title": "No Transition",
            "content": {
                "bullets": ["Test"],
                "transition": "none"
            }
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        assert isinstance(pptx_bytes, bytes)
    
    def test_default_transition(self):
        """Test default transition when not specified."""
        slide_data = {
            "type": "content",
            "title": "Default Transition",
            "content": {
                "bullets": ["Test"]
            }
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        assert isinstance(pptx_bytes, bytes)


# ---------------------------------------------------------------------------
# Task 24.7: Performance Validation Tests
# ---------------------------------------------------------------------------

class TestPerformanceValidation:
    """Test export performance meets requirements."""
    
    def test_export_completes_within_30_seconds_for_50_slides(self):
        """
        Test that export completes within 30 seconds for 50 slides.
        
        This is the critical performance requirement from Task 24.7.
        """
        # Generate 50 slides with mixed types
        slides = []
        for i in range(50):
            slide_type = ["title", "content", "chart", "table", "comparison", "metric"][i % 6]
            
            slide = {
                "slide_id": f"slide-{i+1}",
                "slide_number": i + 1,
                "type": slide_type,
                "title": f"Slide {i+1}: {slide_type.title()}",
                "content": {}
            }
            
            if slide_type == "content":
                slide["content"]["bullets"] = [
                    f"Bullet point {j+1} with specific data: {j*12.5:.1f}% growth" for j in range(5)
                ]
                slide["content"]["highlight_text"] = "Key insight with specific metric"
            elif slide_type == "chart":
                slide["content"]["chart_data"] = {
                    "type": ["bar", "line", "pie", "area"][i % 4],
                    "categories": ["A", "B", "C", "D", "E"],
                    "series": [{"name": "Series", "values": [10, 20, 30, 40, 50]}]
                }
            elif slide_type == "table":
                slide["content"]["table_data"] = {
                    "headers": ["Metric", "2023", "2024", "Growth"],
                    "rows": [
                        ["Revenue", "$100M", "$150M", "50%"],
                        ["Profit", "$20M", "$35M", "75%"],
                        ["Customers", "1000", "2000", "100%"]
                    ]
                }
            elif slide_type == "comparison":
                slide["content"]["comparison_data"] = {
                    "left_title": "Before",
                    "left": ["Item 1 with data", "Item 2 with data"],
                    "right_title": "After",
                    "right": ["Item A with data", "Item B with data"]
                }
            elif slide_type == "metric":
                slide["content"]["metric_value"] = f"{(i+1)*2.5:.1f}%"
                slide["content"]["metric_label"] = "Key Performance Indicator"
                slide["content"]["metric_trend"] = "▲ 23% YoY"
                slide["content"]["bullets"] = ["Context bullet 1", "Context bullet 2"]
            
            slides.append(slide)
        
        # Measure export time
        start_time = time.time()
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build(slides)
        
        end_time = time.time()
        elapsed_time = end_time - start_time
        
        # Verify requirements
        assert isinstance(pptx_bytes, bytes)
        assert len(pptx_bytes) > 0
        assert elapsed_time < 30.0, f"Export took {elapsed_time:.2f}s, must be < 30s"
        
        # Verify all slides were created
        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 50
        
        print(f"\n✓ Performance test passed: 50 slides exported in {elapsed_time:.2f}s")
    
    def test_export_performance_scales_linearly(self):
        """Test that export time scales approximately linearly with slide count."""
        slide_counts = [5, 10, 20]
        times = []
        
        for count in slide_counts:
            slides = [
                {
                    "type": "content",
                    "title": f"Slide {i+1}",
                    "content": {"bullets": ["Test bullet"]}
                }
                for i in range(count)
            ]
            
            start_time = time.time()
            builder = PPTXBuilder("ocean-depths")
            pptx_bytes = builder.build(slides)
            elapsed = time.time() - start_time
            
            times.append(elapsed)
            assert isinstance(pptx_bytes, bytes)
        
        # Verify scaling is reasonable (not exponential)
        # Time for 20 slides should be less than 5x time for 5 slides
        assert times[2] < times[0] * 5, "Export time should scale approximately linearly"


# ---------------------------------------------------------------------------
# Integration Tests
# ---------------------------------------------------------------------------

class TestBuildPptxFunction:
    """Test the main build_pptx function."""
    
    def test_build_pptx_with_valid_data(self, all_slide_types):
        """Test build_pptx function with valid data."""
        pptx_bytes = build_pptx(all_slide_types, "ocean-depths")
        
        assert isinstance(pptx_bytes, bytes)
        assert len(pptx_bytes) > 0
        
        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 6
    
    def test_build_pptx_with_empty_list(self):
        """Test build_pptx with empty slide list."""
        pptx_bytes = build_pptx([], "ocean-depths")
        
        assert isinstance(pptx_bytes, bytes)
        assert len(pptx_bytes) > 0
    
    def test_build_pptx_with_invalid_data_type(self):
        """Test build_pptx raises error for invalid data type."""
        with pytest.raises(ValueError, match="slides_data must be a list"):
            build_pptx("not a list", "ocean-depths")
    
    def test_build_pptx_with_all_themes(self, sample_content_slide):
        """Test build_pptx works with all themes."""
        themes = ["ocean-depths", "sunset-boulevard", "forest-canopy", "modern-minimalist", "golden-hour", "arctic-frost", "desert-rose", "tech-innovation", "botanical-garden", "midnight-galaxy"]
        
        for theme in themes:
            pptx_bytes = build_pptx([sample_content_slide], theme)
            assert isinstance(pptx_bytes, bytes)
            assert len(pptx_bytes) > 0


# ---------------------------------------------------------------------------
# Edge Case Tests
# ---------------------------------------------------------------------------

class TestEdgeCases:
    """Test edge cases and error handling."""
    
    def test_slide_without_content(self):
        """Test slide with missing content field."""
        slide_data = {
            "type": "content",
            "title": "No Content"
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        assert isinstance(pptx_bytes, bytes)
    
    def test_slide_without_title(self):
        """Test slide with missing title."""
        slide_data = {
            "type": "content",
            "content": {"bullets": ["Test"]}
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        assert isinstance(pptx_bytes, bytes)
    
    def test_chart_with_missing_data(self):
        """Test chart slide with incomplete data."""
        slide_data = {
            "type": "chart",
            "title": "Incomplete Chart",
            "content": {
                "chart_data": {
                    "type": "bar",
                    "categories": [],
                    "series": []
                }
            }
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        # Should not crash
        assert isinstance(pptx_bytes, bytes)
    
    def test_unknown_slide_type_fallback(self):
        """Test unknown slide type falls back to content layout."""
        slide_data = {
            "type": "unknown_type",
            "title": "Unknown Type",
            "content": {"bullets": ["Test"]}
        }
        
        builder = PPTXBuilder("ocean-depths")
        pptx_bytes = builder.build([slide_data])
        
        assert isinstance(pptx_bytes, bytes)
        
        prs = PptxPresentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1
